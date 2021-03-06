#!/usr/bin/python
# -*- coding: utf-8 -*-

import os
import yaml
import argparse
from pptx import Presentation
from pptx.util import Cm, Inches
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

DEBUG=False

class PySlides(object):
    level = -1

    def __init__(self, document):
        self.slide_dict = yaml.load(document)

    def generate(self):
        slides = self.slide_dict.get('slides')
        units = self.slide_dict.get('units')
        if units == 'Inches':
            self.unitconv = Inches
        elif units == 'Cm':
            self.unitconv = Cm
        else:
            self.unitconv = Inches
        self.process_slides(slides)
        return self.prs

    def process_slides(self, slides):
        self.prs = Presentation()
        for slide_data in slides:
            self.process_slide(slide_data)

    def process_slide(self, slide_data):
        layout = slide_data.get('layout', 1)
        slide_layout = self.prs.slide_layouts[layout]
        slide = self.prs.slides.add_slide(slide_layout)
        self.shapes = slide.shapes

        for k, v in slide_data.items():
            if DEBUG:
                print( k, v)
            if hasattr(self, k):
                method = getattr(self, k)
                self.process_cmd(method, v)

    def process_cmd(self, method, args):
        if DEBUG:
            print( "process_cmd", method, args, self.level)
        if isinstance(args, list):
            self.level+=1
            for arg in args:
                if DEBUG:
                    print( "Process list item:", arg)
                self.process_cmd(method, arg)
            self.level-=1
        elif isinstance(args, dict):
            if DEBUG: 
                print( "Process dict item", args)
            method(**args)
        else:
            if DEBUG: 
                print( "Process other:", args)
            method(args)

    def title(self, titlestr):
        title = self.shapes.title
        title.text = titlestr

    def text(self, text):
        self.body = self.shapes.placeholders[1]
        tf = self.body.text_frame
        p = tf.add_paragraph()
        p.text = text
        if self.level >= 0: 
            p.level = self.level
        else:
            p.level = 0

    def img(self, path, top, left, width=None, height=None):
        if width and height:
            self.shapes.add_picture(
                path,
                self.unitconv(left),
                self.unitconv(top),
                width=self.unitconv(width),
                height=self.unitconv(height)
            )
        elif width:
            self.shapes.add_picture(
                path,
                self.unitconv(left),
                self.unitconv(top),
                width=self.unitconv(width)
            )
        elif height:
            self.shapes.add_picture(
                path, self.unitconv(left),
                self.unitconv(top),
                height=self.unitconv(height)
            )
        else:
            self.shapes.add_picture(
                path,
                self.unitconv(left),
                self.unitconv(top)
            )

    def chart(self, style, categories, series, x=1, y=2, cx=8, cy=5):
        chart_data = ChartData()
        chart_data.categories = categories
        for data in series:
            chart_data.add_series(data.get('title'), tuple(data.get('data')))

        chart_type = getattr(XL_CHART_TYPE, style)
        chart = self.shapes.add_chart(
            chart_type,
            self.unitconv(x),
            self.unitconv(y),
            self.unitconv(cx),
            self.unitconv(cy),
            chart_data
        ).chart

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.include_in_layout = False


def show_layouts():
    prs = Presentation()
    for i, layout in enumerate(prs.slide_layouts):
        print( i, layout.name)


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument('file', nargs='?', help='yaml file')
    parser.add_argument(
        '-l', '--layouts', help='List layouts ids', action='store_true'
    )

    args = parser.parse_args()

    if args.layouts:
        show_layouts()
        exit()

    file_name = args.file
    if not file_name:
        parser.print(_help())
        exit()

    if file_name.endswith('.yml') or file_name.endswith('.yaml'):

        file_content = open(file_name).read()
        file_name = os.path.splitext(file_name)[0]

        pptx_name = '{}.pptx'.format(file_name)

        ps = PySlides(file_content)
        prs = ps.generate()
        prs.save(pptx_name)


if __name__ == '__main__':
    main()
